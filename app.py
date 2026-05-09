import os
import re
from datetime import datetime

import streamlit as st
from dotenv import load_dotenv
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain_community.vectorstores import FAISS
from langchain_core.documents import Document
from langchain_groq import ChatGroq
from langchain_text_splitters import RecursiveCharacterTextSplitter
from pypdf import PdfReader
from rank_bm25 import BM25Okapi
from sentence_transformers import CrossEncoder

try:
    from docx import Document as DocxReader
except ImportError:
    DocxReader = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None


load_dotenv()
GROQ_API_KEY = os.getenv("GROQ_API_KEY")

st.set_page_config(
    page_title="ResearchMind AI",
    page_icon="🧠",
    layout="wide"
)


st.markdown(
    """
    <style>
    .main-title {
        font-size: 42px;
        font-weight: 800;
        color: #4A90E2;
        margin-bottom: 0px;
    }

    .subtitle {
        color: gray;
        font-size: 17px;
        margin-top: 0px;
    }
    </style>
    """,
    unsafe_allow_html=True
)

st.markdown(
    '<p class="main-title">🧠 ResearchMind AI</p>',
    unsafe_allow_html=True
)

st.markdown(
    '<p class="subtitle">Advanced AI Research Assistant with PDF/DOCX/PPTX support, citations, hybrid search, reranking, and answer evaluation</p>',
    unsafe_allow_html=True
)


if "messages" not in st.session_state:
    st.session_state.messages = []

if "vector_store" not in st.session_state:
    st.session_state.vector_store = None

if "chunks" not in st.session_state:
    st.session_state.chunks = []

if "raw_text" not in st.session_state:
    st.session_state.raw_text = ""

if "processed_file_names" not in st.session_state:
    st.session_state.processed_file_names = []

if "bm25" not in st.session_state:
    st.session_state.bm25 = None

if "document_summary" not in st.session_state:
    st.session_state.document_summary = ""

if "flashcards" not in st.session_state:
    st.session_state.flashcards = ""

if "mcqs" not in st.session_state:
    st.session_state.mcqs = ""

if "last_question" not in st.session_state:
    st.session_state.last_question = ""

if "last_answer" not in st.session_state:
    st.session_state.last_answer = ""

if "last_context" not in st.session_state:
    st.session_state.last_context = ""

if "answer_evaluation" not in st.session_state:
    st.session_state.answer_evaluation = ""


@st.cache_resource(show_spinner=False)
def load_embeddings():
    return HuggingFaceEmbeddings(
        model_name="sentence-transformers/all-MiniLM-L6-v2"
    )


@st.cache_resource(show_spinner=False)
def load_reranker():
    return CrossEncoder(
        "cross-encoder/ms-marco-MiniLM-L-6-v2"
    )


def get_llm(model_name, temperature):
    return ChatGroq(
        groq_api_key=GROQ_API_KEY,
        model_name=model_name,
        temperature=temperature
    )


def run_llm(prompt, model_name, temperature):
    llm = get_llm(model_name, temperature)
    response = llm.invoke(prompt)
    return response.content


def simple_tokenize(text):
    return re.findall(r"\b\w+\b", text.lower())


def extract_pdf(uploaded_file):
    documents = []
    full_text = ""
    file_name = uploaded_file.name
    pdf_reader = PdfReader(uploaded_file)

    for page_number, page in enumerate(pdf_reader.pages, start=1):
        extracted_text = page.extract_text()

        if extracted_text:
            metadata = {
                "source": file_name,
                "type": "PDF",
                "location": f"Page {page_number}"
            }

            documents.append(
                Document(
                    page_content=extracted_text,
                    metadata=metadata
                )
            )

            full_text += f"\n\n[Source: {file_name}, Page {page_number}]\n"
            full_text += extracted_text

    return documents, full_text


def extract_docx(uploaded_file):
    documents = []
    full_text = ""
    file_name = uploaded_file.name

    if DocxReader is None:
        st.error("DOCX support requires python-docx. Run: pip install python-docx")
        return documents, full_text

    doc = DocxReader(uploaded_file)
    paragraphs = []

    for paragraph in doc.paragraphs:
        if paragraph.text.strip():
            paragraphs.append(paragraph.text.strip())

    text = "\n".join(paragraphs)

    if text:
        metadata = {
            "source": file_name,
            "type": "DOCX",
            "location": "Document"
        }

        documents.append(
            Document(
                page_content=text,
                metadata=metadata
            )
        )

        full_text += f"\n\n[Source: {file_name}, DOCX Document]\n"
        full_text += text

    return documents, full_text


def extract_pptx(uploaded_file):
    documents = []
    full_text = ""
    file_name = uploaded_file.name

    if Presentation is None:
        st.error("PPTX support requires python-pptx. Run: pip install python-pptx")
        return documents, full_text

    presentation = Presentation(uploaded_file)

    for slide_number, slide in enumerate(presentation.slides, start=1):
        slide_text_parts = []

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text_parts.append(shape.text.strip())

        slide_text = "\n".join(slide_text_parts)

        if slide_text:
            metadata = {
                "source": file_name,
                "type": "PPTX",
                "location": f"Slide {slide_number}"
            }

            documents.append(
                Document(
                    page_content=slide_text,
                    metadata=metadata
                )
            )

            full_text += f"\n\n[Source: {file_name}, Slide {slide_number}]\n"
            full_text += slide_text

    return documents, full_text


def extract_text_from_files(uploaded_files):
    all_documents = []
    complete_text = ""

    for uploaded_file in uploaded_files:
        file_name = uploaded_file.name.lower()

        if file_name.endswith(".pdf"):
            documents, text = extract_pdf(uploaded_file)

        elif file_name.endswith(".docx"):
            documents, text = extract_docx(uploaded_file)

        elif file_name.endswith(".pptx"):
            documents, text = extract_pptx(uploaded_file)

        else:
            documents, text = [], ""

        all_documents.extend(documents)
        complete_text += text

    return all_documents, complete_text


def build_vector_store(documents):
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=1500,
        chunk_overlap=300,
        separators=["\n\n", "\n", ".", " "]
    )

    chunks = text_splitter.split_documents(documents)

    for index, chunk in enumerate(chunks):
        chunk.metadata["chunk_id"] = index

    embeddings = load_embeddings()

    vector_store = FAISS.from_documents(
        chunks,
        embedding=embeddings
    )

    tokenized_chunks = []

    for chunk in chunks:
        tokenized_chunks.append(
            simple_tokenize(chunk.page_content)
        )

    bm25 = BM25Okapi(tokenized_chunks)

    return vector_store, chunks, bm25


def format_source_label(doc):
    source = doc.metadata.get("source", "Unknown source")
    location = doc.metadata.get("location", "Unknown location")
    return f"{source} — {location}"


def create_context_with_citations(docs):
    context_parts = []

    for index, doc in enumerate(docs, start=1):
        label = format_source_label(doc)

        context_parts.append(
            f"[Source {index}: {label}]\n{doc.page_content}"
        )

    return "\n\n".join(context_parts)


def hybrid_retrieve(
    query,
    vector_store,
    chunks,
    bm25,
    semantic_k,
    keyword_k,
    final_k,
    use_reranker
):
    semantic_docs = vector_store.similarity_search(
        query,
        k=semantic_k
    )

    query_tokens = simple_tokenize(query)
    bm25_scores = bm25.get_scores(query_tokens)

    top_keyword_indices = sorted(
        range(len(bm25_scores)),
        key=lambda index: bm25_scores[index],
        reverse=True
    )[:keyword_k]

    keyword_docs = []

    for index in top_keyword_indices:
        keyword_docs.append(chunks[index])

    combined_docs = []
    seen_ids = set()

    for doc in semantic_docs + keyword_docs:
        chunk_id = doc.metadata.get("chunk_id")

        if chunk_id not in seen_ids:
            combined_docs.append(doc)
            seen_ids.add(chunk_id)

    if use_reranker and combined_docs:
        reranker = load_reranker()

        pairs = []

        for doc in combined_docs:
            pairs.append((query, doc.page_content))

        scores = reranker.predict(pairs)

        ranked_pairs = sorted(
            zip(combined_docs, scores),
            key=lambda item: item[1],
            reverse=True
        )

        reranked_docs = []

        for doc, score in ranked_pairs[:final_k]:
            reranked_docs.append(doc)

        return reranked_docs

    return combined_docs[:final_k]


def create_download_text():
    conversation = "ResearchMind AI Chat Export\n"
    conversation += f"Exported on: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n\n"

    for message in st.session_state.messages:
        role = message["role"].upper()
        content = message["content"]
        conversation += f"{role}:\n{content}\n\n"

    return conversation


with st.sidebar:
    st.header("⚙️ Settings")

    model_name = st.selectbox(
        "Choose Groq Model",
        [
            "llama-3.3-70b-versatile",
            "deepseek-r1-distill-llama-70b",
            "mixtral-8x7b-32768",
            "llama-3.1-8b-instant"
        ],
        index=0
    )

    temperature = st.slider(
        "Creativity / Temperature",
        min_value=0.0,
        max_value=1.0,
        value=0.2,
        step=0.1
    )

    st.divider()
    st.subheader("🔍 Retrieval Settings")

    semantic_k = st.slider(
        "Semantic Search Results",
        min_value=1,
        max_value=10,
        value=6
    )

    keyword_k = st.slider(
        "Keyword Search Results",
        min_value=1,
        max_value=10,
        value=6
    )

    final_k = st.slider(
        "Final Source Chunks",
        min_value=1,
        max_value=8,
        value=5
    )

    use_reranker = st.checkbox(
        "Use Advanced Reranking",
        value=True
    )

    st.divider()

    uploaded_files = st.file_uploader(
        "📄 Upload PDF, DOCX, or PPTX Files",
        type=["pdf", "docx", "pptx"],
        accept_multiple_files=True
    )

    process_button = st.button("⚙️ Process Uploaded Files")

    if st.button("🧹 Clear Chat"):
        st.session_state.messages = []
        st.rerun()

    if st.button("🔄 Reset Documents"):
        st.session_state.vector_store = None
        st.session_state.chunks = []
        st.session_state.raw_text = ""
        st.session_state.processed_file_names = []
        st.session_state.bm25 = None
        st.session_state.document_summary = ""
        st.session_state.flashcards = ""
        st.session_state.mcqs = ""
        st.session_state.answer_evaluation = ""
        st.rerun()


if not GROQ_API_KEY:
    st.error("Missing GROQ_API_KEY. Please add it to your .env file.")
    st.code("GROQ_API_KEY=your_groq_api_key_here")
    st.stop()


if uploaded_files and process_button:
    current_file_names = []

    for uploaded_file in uploaded_files:
        current_file_names.append(uploaded_file.name)

    with st.spinner("Extracting text and building vector database..."):
        documents, raw_text = extract_text_from_files(uploaded_files)

        if not documents:
            st.error("No readable text found in the uploaded files.")

        else:
            vector_store, chunks, bm25 = build_vector_store(documents)

            st.session_state.raw_text = raw_text
            st.session_state.vector_store = vector_store
            st.session_state.chunks = chunks
            st.session_state.bm25 = bm25
            st.session_state.processed_file_names = current_file_names

            st.success("✅ Files processed successfully!")


col1, col2, col3, col4 = st.columns(4)

with col1:
    st.metric("Documents", len(st.session_state.processed_file_names))

with col2:
    st.metric("Text Chunks", len(st.session_state.chunks))

with col3:
    st.metric("Chat Messages", len(st.session_state.messages))

with col4:
    if use_reranker:
        st.metric("Retrieval", "Hybrid + Rerank")
    else:
        st.metric("Retrieval", "Hybrid")


st.divider()
st.subheader("🛠️ Research Tools")

tool_col1, tool_col2, tool_col3, tool_col4 = st.columns(4)

with tool_col1:
    if st.button("📘 Generate Summary"):
        if not st.session_state.raw_text:
            st.warning("Please upload and process a document first.")

        else:
            with st.spinner("Generating document summary..."):
                prompt = f"""
You are an expert academic research assistant.

Create a university-level summary of the document.

Include:
1. Main topic
2. Key ideas
3. Important findings
4. Technical terms explained simply
5. Short exam/revision notes

DOCUMENT:
{st.session_state.raw_text[:18000]}
"""

                st.session_state.document_summary = run_llm(
                    prompt,
                    model_name,
                    temperature
                )

with tool_col2:
    if st.button("🧠 Generate Flashcards"):
        if not st.session_state.raw_text:
            st.warning("Please upload and process a document first.")

        else:
            with st.spinner("Generating flashcards..."):
                prompt = f"""
Create high-quality study flashcards from this document.

Format exactly like this:
Q1: ...
A1: ...

Make them useful for revision.

DOCUMENT:
{st.session_state.raw_text[:15000]}
"""

                st.session_state.flashcards = run_llm(
                    prompt,
                    model_name,
                    temperature
                )

with tool_col3:
    if st.button("📝 Generate MCQs"):
        if not st.session_state.raw_text:
            st.warning("Please upload and process a document first.")

        else:
            with st.spinner("Generating MCQs..."):
                prompt = f"""
Generate 15 exam-style multiple choice questions from this document.

For each question include:
- Question
- Four options: A, B, C, D
- Correct answer
- Short explanation

DOCUMENT:
{st.session_state.raw_text[:15000]}
"""

                st.session_state.mcqs = run_llm(
                    prompt,
                    model_name,
                    temperature
                )

with tool_col4:
    if st.button("✅ Evaluate Last Answer"):
        if not st.session_state.last_answer:
            st.warning("Ask a question first before evaluating an answer.")

        else:
            with st.spinner("Evaluating answer quality..."):
                prompt = f"""
You are an answer quality evaluator for a RAG system.

Evaluate whether the answer is supported by the provided context.

Return:
1. Groundedness score out of 10
2. Completeness score out of 10
3. Accuracy concerns
4. Missing information
5. How to improve the answer

QUESTION:
{st.session_state.last_question}

CONTEXT:
{st.session_state.last_context}

ANSWER:
{st.session_state.last_answer}
"""

                st.session_state.answer_evaluation = run_llm(
                    prompt,
                    model_name,
                    0.0
                )


if st.session_state.document_summary:
    with st.expander("📘 Document Summary", expanded=False):
        st.markdown(st.session_state.document_summary)

if st.session_state.flashcards:
    with st.expander("🧠 Flashcards", expanded=False):
        st.markdown(st.session_state.flashcards)

if st.session_state.mcqs:
    with st.expander("📝 MCQs", expanded=False):
        st.markdown(st.session_state.mcqs)

if st.session_state.answer_evaluation:
    with st.expander("✅ Answer Evaluation", expanded=True):
        st.markdown(st.session_state.answer_evaluation)


st.divider()
st.subheader("💬 Chat with Your Documents")

for message in st.session_state.messages:
    with st.chat_message(message["role"]):
        st.markdown(message["content"])


user_question = st.chat_input("Ask a question from your documents...")

if user_question:
    st.session_state.messages.append(
        {
            "role": "user",
            "content": user_question
        }
    )

    with st.chat_message("user"):
        st.markdown(user_question)

    if st.session_state.vector_store is None or st.session_state.bm25 is None:
        with st.chat_message("assistant"):
            warning_message = "⚠️ Please upload and process at least one document before asking questions."
            st.warning(warning_message)

            st.session_state.messages.append(
                {
                    "role": "assistant",
                    "content": warning_message
                }
            )

    else:
        with st.chat_message("assistant"):
            with st.spinner("Searching sources and generating answer..."):
                try:
                    docs = hybrid_retrieve(
                        query=user_question,
                        vector_store=st.session_state.vector_store,
                        chunks=st.session_state.chunks,
                        bm25=st.session_state.bm25,
                        semantic_k=semantic_k,
                        keyword_k=keyword_k,
                        final_k=final_k,
                        use_reranker=use_reranker
                    )

                    context = create_context_with_citations(docs)

                    prompt = f"""
You are ResearchMind AI, an advanced academic research assistant.

Rules:
- Answer ONLY using the provided document context.
- Use citations like [Source 1], [Source 2] after important claims.
- If the answer is not present, say: "I could not find this information in the uploaded documents."
- Explain clearly and logically.
- Use bullet points where helpful.
- Mention important terms and explain them simply.
- Avoid hallucination.

DOCUMENT CONTEXT:
{context}

USER QUESTION:
{user_question}

FINAL ANSWER WITH CITATIONS:
"""

                    answer = run_llm(
                        prompt,
                        model_name,
                        temperature
                    )

                    st.markdown(answer)

                    st.session_state.last_question = user_question
                    st.session_state.last_answer = answer
                    st.session_state.last_context = context

                    st.session_state.messages.append(
                        {
                            "role": "assistant",
                            "content": answer
                        }
                    )

                    with st.expander("📚 Source Citations Used"):
                        for i, doc in enumerate(docs, start=1):
                            st.markdown(
                                f"### Source {i}: {format_source_label(doc)}"
                            )

                            st.write(
                                doc.page_content[:2000]
                            )

                except Exception as e:
                    error_message = f"Error: {e}"
                    st.error(error_message)

                    st.session_state.messages.append(
                        {
                            "role": "assistant",
                            "content": error_message
                        }
                    )


st.divider()

if st.session_state.messages:
    export_text = create_download_text()

    st.download_button(
        label="💾 Download Chat History",
        data=export_text,
        file_name=f"researchmind_chat_{datetime.now().strftime('%Y%m%d_%H%M%S')}.txt",
        mime="text/plain"
    )

st.caption(
    "🚀 Powered by Streamlit + LangChain + FAISS + BM25 + Cross-Encoder Reranking + Groq"
)